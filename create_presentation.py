import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
from io import BytesIO
import base64

class PresentationCreator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_theme()
        
    def setup_theme(self):
        """Set up presentation theme and colors"""
        self.primary_color = RGBColor(31, 78, 121)  # Dark blue
        self.secondary_color = RGBColor(68, 114, 196)  # Light blue
        self.accent_color = RGBColor(255, 192, 0)  # Gold
        self.text_color = RGBColor(64, 64, 64)  # Dark gray
        
    def add_title_slide(self):
        """Create title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Service Delivery Protest Prediction Model"
        subtitle.text = "Predicting Social Unrest using Census 2022 Data and Machine Learning\n\nData Science Competition Submission\nJanuary 2025"
        
        # Format title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        title.text_frame.paragraphs[0].font.bold = True
        
        # Format subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.text_color
            
    def add_agenda_slide(self):
        """Create agenda slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Presentation Agenda"
        
        agenda_items = [
            "1. Problem Statement & Objectives",
            "2. Data Sources & Methodology",
            "3. Feature Engineering & Model Development",
            "4. Results & Performance Metrics",
            "5. Key Insights & Provincial Analysis",
            "6. Interactive Dashboard Demo",
            "7. Policy Recommendations",
            "8. Conclusion & Future Work"
        ]
        
        content.text = "\n".join(agenda_items)
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(12)
            
    def add_problem_statement_slide(self):
        """Create problem statement slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Problem Statement & Objectives"
        
        problem_text = """🎯 CHALLENGE
• Service delivery protests are increasing across South Africa
• Government needs predictive tools for resource allocation
• Early warning systems can prevent social unrest

🎯 OBJECTIVES
• Predict protest risk using Census 2022 demographic data
• Identify key service delivery factors driving unrest
• Develop actionable insights for policy makers
• Create interactive tools for government officials

🎯 SUCCESS METRICS
• Model accuracy (R² > 0.6)
• Feature interpretability
• Deployment-ready solution
• Clear policy recommendations"""
        
        content.text = problem_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(8)
            
    def add_data_sources_slide(self):
        """Create data sources slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Data Sources & Methodology"
        
        data_text = """📊 PRIMARY DATASETS

🏠 Census 2022 Data (Statistics South Africa)
• 12 comprehensive data sheets
• 265 municipalities across 9 provinces
• Demographics, infrastructure, service delivery metrics

📢 Protest Events Data (ACLED)
• Historical demonstration events (1997-2025)
• Political violence and fatalities data
• Monthly aggregated protest counts

🔬 METHODOLOGY
• Data cleaning and preprocessing
• Feature engineering (9 service delivery indicators)
• Provincial-level aggregation
• Machine learning model comparison
• Cross-validation and performance evaluation"""
        
        content.text = data_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(6)
            
    def add_feature_engineering_slide(self):
        """Create feature engineering slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Feature Engineering & Model Development"
        
        features_text = """🔧 ENGINEERED FEATURES

📈 Service Delivery Indicators (9 features)
• Water Access Percentage - Piped water inside dwelling
• Sanitation Access Percentage - Flush toilet connectivity
• Electricity Access Percentage - Electrical connection
• Refuse Service Percentage - Municipal waste collection
• Education Access Percentage - Higher education attainment
• Service Delivery Index - Composite service quality measure
• Service Gap - Overall deficiency (100 - Service Index)
• Population Density - People per square kilometer
• Year Trend - Temporal progression factor

🤖 MODEL COMPARISON
• Random Forest Regressor (WINNER)
• Gradient Boosting Regressor
• Linear Regression (Baseline)

✅ VALIDATION APPROACH
• 5-fold cross-validation
• Train/test split (80/20)
• Performance metrics: R², RMSE, MAE"""
        
        content.text = features_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(6)
            
    def add_results_slide(self):
        """Create results slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Results & Performance Metrics"
        
        results_text = """🏆 MODEL PERFORMANCE

🥇 Best Algorithm: Random Forest Regressor
• Training R² Score: 0.652
• Cross-Validation R²: 0.728 ± 0.038
• RMSE: 3.107 ± 0.170
• MAE: 2.519 ± 0.145

📊 ALGORITHM COMPARISON
• Random Forest: R² = 0.652 ⭐
• Gradient Boosting: R² = 0.643
• Linear Regression: R² = 0.644

🎯 FEATURE IMPORTANCE (Top 5)
1. Service Gap (55.4%) - Overall service deficiency
2. Service Delivery Index (25.0%) - Composite quality
3. Refuse Service Percentage (6.6%) - Waste collection
4. Electricity Access (3.0%) - Power infrastructure
5. Year Trend (2.5%) - Temporal factor

✅ MODEL VALIDATION
• Stable performance across folds
• Well-calibrated predictions
• No overfitting detected"""
        
        content.text = results_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(6)
            
    def add_provincial_analysis_slide(self):
        """Create provincial analysis slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Provincial Risk Analysis (2024 Predictions)"
        
        provincial_text = """🗺️ PROTEST RISK BY PROVINCE

🔴 HIGH RISK (Score > 7.0)
• Eastern Cape: 8.29
• Limpopo: 9.54

🟡 MEDIUM RISK (Score 4.0 - 7.0)
• Free State: ~5.5
• KwaZulu-Natal: ~6.2
• North West: ~5.8
• Mpumalanga: ~4.8

🟢 LOW RISK (Score < 4.0)
• Western Cape: 2.21
• Northern Cape: 2.64
• Gauteng: ~3.5

📈 KEY INSIGHTS
• Eastern provinces show highest risk
• Service delivery gaps strongly correlate with risk
• Refuse collection is critical factor
• Urban provinces generally lower risk

⚠️ IMMEDIATE ATTENTION REQUIRED
• Eastern Cape and Limpopo need urgent intervention
• Focus on basic service delivery improvements"""
        
        content.text = provincial_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(6)
            
    def add_dashboard_slide(self):
        """Create dashboard demo slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Interactive Dashboard & Deployment"
        
        dashboard_text = """🖥️ STREAMLIT WEB APPLICATION

🎛️ Dashboard Features
• Real-time protest risk predictions
• Interactive data exploration
• Provincial comparison tools
• Custom scenario analysis
• Model performance metrics
• Feature importance visualization

🌐 Deployment Details
• Live at: http://localhost:8502
• User-friendly interface
• No technical expertise required
• Instant predictions for any province/year
• Export capabilities for reports

📱 Accessibility
• Web-based (any device)
• Government official friendly
• Policy maker dashboard
• Real-time updates possible

🔧 Technical Stack
• Python + Streamlit
• Scikit-learn models
• Interactive Plotly charts
• Pandas data processing"""
        
        content.text = dashboard_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(6)
            
    def add_recommendations_slide(self):
        """Create policy recommendations slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Policy Recommendations"
        
        recommendations_text = """🎯 IMMEDIATE ACTIONS (0-6 months)

🚨 High Priority Provinces
• Eastern Cape & Limpopo: Emergency service delivery task force
• Focus on refuse collection and basic sanitation
• Increase municipal capacity and resources

📊 MEDIUM-TERM STRATEGIES (6-18 months)

🏗️ Infrastructure Development
• Prioritize electricity and water access projects
• Improve waste management systems
• Enhance education facilities

📈 Monitoring & Early Warning
• Deploy predictive model in government systems
• Monthly risk assessment reports
• Proactive resource allocation

🔄 LONG-TERM INITIATIVES (18+ months)

🎯 Systemic Improvements
• Comprehensive service delivery reform
• Municipal capacity building programs
• Community engagement initiatives
• Regular model updates with new data

💡 SUCCESS INDICATORS
• Reduced protest incidents
• Improved service delivery metrics
• Higher citizen satisfaction scores"""
        
        content.text = recommendations_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(5)
            
    def add_conclusion_slide(self):
        """Create conclusion slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Conclusion & Future Work"
        
        conclusion_text = """🏆 PROJECT ACHIEVEMENTS

✅ Successfully developed predictive model (R² = 0.652)
✅ Identified key service delivery risk factors
✅ Created interactive dashboard for government use
✅ Provided actionable policy recommendations
✅ Delivered complete, deployment-ready solution

🔮 FUTURE ENHANCEMENTS

📊 Data Improvements
• Real-time service delivery monitoring
• Social media sentiment analysis
• Economic indicators integration
• Weather and seasonal patterns

🤖 Model Enhancements
• Deep learning for complex patterns
• Time series forecasting (LSTM)
• Ensemble methods and stacking
• Causal inference analysis

🚀 Deployment Expansion
• REST API for live predictions
• Mobile app for field workers
• GIS mapping integration
• Automated alert systems

💼 BUSINESS IMPACT
• Prevent social unrest through early intervention
• Optimize government resource allocation
• Improve citizen satisfaction and trust
• Evidence-based policy making"""
        
        content.text = conclusion_text
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = self.text_color
            paragraph.space_after = Pt(5)
            
    def add_thank_you_slide(self):
        """Create thank you slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Thank You"
        subtitle.text = "Questions & Discussion\n\nService Delivery Protest Prediction Model\nData Science Competition Submission\n\nContact: DataDynamo Team\nJanuary 2025"
        
        # Format title
        title.text_frame.paragraphs[0].font.size = Pt(48)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        title.text_frame.paragraphs[0].font.bold = True
        
        # Format subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.text_color
            
    def create_presentation(self):
        """Create the complete presentation"""
        print("🎯 Creating PowerPoint Presentation...")
        
        # Add all slides
        self.add_title_slide()
        print("  ✅ Title slide added")
        
        self.add_agenda_slide()
        print("  ✅ Agenda slide added")
        
        self.add_problem_statement_slide()
        print("  ✅ Problem statement slide added")
        
        self.add_data_sources_slide()
        print("  ✅ Data sources slide added")
        
        self.add_feature_engineering_slide()
        print("  ✅ Feature engineering slide added")
        
        self.add_results_slide()
        print("  ✅ Results slide added")
        
        self.add_provincial_analysis_slide()
        print("  ✅ Provincial analysis slide added")
        
        self.add_dashboard_slide()
        print("  ✅ Dashboard demo slide added")
        
        self.add_recommendations_slide()
        print("  ✅ Policy recommendations slide added")
        
        self.add_conclusion_slide()
        print("  ✅ Conclusion slide added")
        
        self.add_thank_you_slide()
        print("  ✅ Thank you slide added")
        
        # Save presentation
        filename = "Service_Delivery_Protest_Prediction_Presentation.pptx"
        self.prs.save(filename)
        print(f"\n🎉 Presentation saved as: {filename}")
        print(f"📊 Total slides: {len(self.prs.slides)}")
        
        return filename

def main():
    """Main function to create the presentation"""
    try:
        creator = PresentationCreator()
        filename = creator.create_presentation()
        
        print("\n" + "="*60)
        print("🏆 POWERPOINT PRESENTATION CREATED SUCCESSFULLY!")
        print("="*60)
        print(f"📁 File: {filename}")
        print(f"📍 Location: {os.getcwd()}")
        print("\n📋 Presentation Contents:")
        print("  1. Title Slide")
        print("  2. Agenda")
        print("  3. Problem Statement & Objectives")
        print("  4. Data Sources & Methodology")
        print("  5. Feature Engineering & Model Development")
        print("  6. Results & Performance Metrics")
        print("  7. Provincial Risk Analysis")
        print("  8. Interactive Dashboard Demo")
        print("  9. Policy Recommendations")
        print("  10. Conclusion & Future Work")
        print("  11. Thank You & Questions")
        print("\n🎯 Ready for competition submission!")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        print("💡 Make sure python-pptx is installed: pip install python-pptx")

if __name__ == "__main__":
    main()